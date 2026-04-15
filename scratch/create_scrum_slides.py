from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def update_scrum_slides():
    prs = Presentation('Architecture.pptx')
    
    # Remove existing slides after Architecture
    while len(prs.slides) > 1:
        rId = prs.slides._sldIdLst[-1].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[-1]

    blank_slide_layout = prs.slide_layouts[6]
    
    BLUE_DARK = RGBColor(13, 33, 55)
    BLUE_AZURE = RGBColor(0, 120, 212)
    BLUE_LIGHT = RGBColor(227, 242, 253)
    GREEN = RGBColor(67, 160, 71)
    ORANGE = RGBColor(255, 109, 0)
    GRAY_LIGHT = RGBColor(230, 230, 230)
    
    # --- SLIDE 2: THE SCRUM SPIRAL (ITERATIVE RELEASE CHAIN) ---
    slide2 = prs.slides.add_slide(blank_slide_layout)
    title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(12.33), Inches(0.6))
    p1 = title2.text_frame.paragraphs[0]
    p1.text = "The Scrum Spiral: Iterative Sprints within a Release Cycle"
    p1.font.bold = True
    p1.font.size = Pt(22)
    p1.font.color.rgb = BLUE_DARK

    # 1. Release Container
    boundary = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.2), Inches(1.0), Inches(12.93), Inches(5.8))
    boundary.fill.solid()
    boundary.fill.fore_color.rgb = RGBColor(255, 255, 255)
    boundary.line.color.rgb = GRAY_LIGHT
    boundary.line.dash_style = 2
    
    b_label = slide2.shapes.add_textbox(Inches(0.3), Inches(1.1), Inches(4), Inches(0.4))
    b_label.text_frame.paragraphs[0].text = "RELEASE 1.0: MVP SIMULATION"
    b_label.text_frame.paragraphs[0].font.size = Pt(10)
    b_label.text_frame.paragraphs[0].font.bold = True
    b_label.text_frame.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)

    # 2. Product Backlog
    sh_backlog = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(2.5), Inches(1.2), Inches(2.5))
    sh_backlog.fill.solid()
    sh_backlog.fill.fore_color.rgb = BLUE_LIGHT
    sh_backlog.line.color.rgb = BLUE_AZURE
    sh_backlog.text_frame.paragraphs[0].text = "PRODUCT\nBACKLOG"
    sh_backlog.text_frame.paragraphs[0].font.size = Pt(9)
    sh_backlog.text_frame.paragraphs[0].font.bold = True
    sh_backlog.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 3. THE SPIRAL (Loops)
    def add_sprint_loop(x_off, num):
        arc = slide2.shapes.add_shape(MSO_SHAPE.BLOCK_ARC, Inches(x_off), Inches(2.2), Inches(2.8), Inches(2.8))
        arc.fill.solid()
        arc.fill.fore_color.rgb = BLUE_AZURE
        arc.rotation = 90
        
        lbl = slide2.shapes.add_textbox(Inches(x_off + 0.4), Inches(3.2), Inches(2.0), Inches(0.8))
        p = lbl.text_frame.paragraphs[0]
        p.text = f"SPRINT {num}"
        p.font.bold = True
        p.font.size = Pt(12)
        p.alignment = PP_ALIGN.CENTER
        
        rit = slide2.shapes.add_textbox(Inches(x_off + 0.5), Inches(5.1), Inches(1.8), Inches(0.8))
        p2 = rit.text_frame.paragraphs[0]
        p2.text = "- Planning\n- Daily Standup\n- Review/Retro"
        p2.font.size = Pt(8)
        p2.alignment = PP_ALIGN.CENTER

    add_sprint_loop(2.0, 1)
    add_sprint_loop(4.8, 2)
    add_sprint_loop(7.6, 3)

    # 4. Growing Increment (Chevrons)
    for i in range(3):
        w = 1.0 + (i * 0.8)
        inc = slide2.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(2.8 + (i * 2.8)), Inches(4.5), Inches(w), Inches(0.4))
        inc.fill.solid()
        inc.fill.fore_color.rgb = GREEN
        inc.line.color.rgb = BLUE_DARK
        
    # 5. Production Release
    rel = slide2.shapes.add_shape(MSO_SHAPE.PENTAGON, Inches(10.8), Inches(2.2), Inches(2.0), Inches(2.8))
    rel.fill.solid()
    rel.fill.fore_color.rgb = BLUE_DARK
    rel.text_frame.paragraphs[0].text = "PRODUCTION\nRELEASE"
    rel.text_frame.paragraphs[0].font.bold = True
    rel.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    rel.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    conn = slide2.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(1.6), Inches(3.4), Inches(0.6), Inches(0.4))
    conn.fill.solid()
    conn.fill.fore_color.rgb = BLUE_AZURE

    # --- SLIDE 3: GOVERNANCE ---
    slide3 = prs.slides.add_slide(blank_slide_layout)
    title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.33), Inches(0.8))
    title3.text_frame.paragraphs[0].text = "Project Governance: Managing Constraints"
    title3.text_frame.paragraphs[0].font.bold = True
    title3.text_frame.paragraphs[0].font.size = Pt(24)
    title3.text_frame.paragraphs[0].font.color.rgb = BLUE_DARK

    def add_metric(x, y, title, content):
        box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.8), Inches(2))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        box.line.color.rgb = BLUE_AZURE
        tf = slide3.shapes.add_textbox(Inches(x+0.2), Inches(y+0.2), Inches(5.4), Inches(1.6)).text_frame
        p_t = tf.paragraphs[0]
        p_t.text = title.upper()
        p_t.font.bold = True
        p_t.font.size = Pt(12)
        p_t.font.color.rgb = BLUE_AZURE
        p_c = tf.add_paragraph()
        p_c.text = content
        p_c.font.size = Pt(10)

    add_metric(0.5, 1.5, "Risk Management", "- Empirical Process Control: Transparency, Inspection, and Adaptation.\n- Early Failure Detection: Technical risks identified during 'Sprint 0'.\n- Regular Stakeholder Feedback: Bi-weekly demos with JA Team.")
    add_metric(6.8, 1.5, "Schedule Control", "- Time-Boxing: Fixed durations (12 months) with hard Sprint deadlines.\n- Burndown Charts: Real-time tracking of development velocity.\n- Predictable Cadence: Consistent release cycles for testing.")
    add_metric(0.5, 4.0, "Cost Management", "- $1M Budget Adherence: Agile prioritization ensures high-value features first.\n- No Waste: Focus on validated requirements prevents over-engineering.\n- Resource Efficiency: Cross-functional team minimizes delays.")
    add_metric(6.8, 4.0, "Time Control", "- Sprints: 2-week focus cycles prevent scope creep.\n- Daily Scrum: 15-min resets ensure blockers are removed immediately.\n- Definition of Done: Prevents '99% complete' syndrome.")

    prs.save('Architecture.pptx')
    print("Architecture.pptx finalized with Scrum Spiral.")

if __name__ == "__main__":
    update_scrum_slides()
