from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # Set slide size to Widescreen (16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Colors
    BLUE_DARK = RGBColor(13, 33, 55)
    BLUE_AZURE = RGBColor(0, 120, 212)
    BLUE_LIGHT = RGBColor(227, 242, 253)
    ORANGE = RGBColor(255, 109, 0)
    ORANGE_LIGHT = RGBColor(255, 243, 224)
    PURPLE = RGBColor(123, 31, 162)
    PURPLE_LIGHT = RGBColor(237, 231, 246)
    RED = RGBColor(184, 84, 80)
    RED_LIGHT = RGBColor(255, 235, 238)
    GREEN = RGBColor(67, 160, 71)
    GREEN_LIGHT = RGBColor(241, 248, 233)
    YELLOW = RGBColor(249, 168, 37)
    YELLOW_LIGHT = RGBColor(255, 249, 196)

    def add_layer_bar(y, height, title, color_bg, color_border):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.2), Inches(y), Inches(12.93), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color_bg
        shape.line.color.rgb = color_border
        
        # Title
        text_frame = slide.shapes.add_textbox(Inches(0.3), Inches(y + 0.05), Inches(5), Inches(0.3)).text_frame
        p = text_frame.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(10)
        p.font.color.rgb = BLUE_DARK
        return shape

    def add_box(x, y, title, subtitle, color_border, bg_color=None):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(1.4), Inches(0.9))
        shape.fill.solid()
        if bg_color:
            shape.fill.fore_color.rgb = bg_color
        else:
            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
        shape.line.color.rgb = color_border
        
        # Title
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y + 0.1), Inches(1.4), Inches(0.3))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(9)
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(7)
        p2.font.color.rgb = RGBColor(80, 80, 80)
        p2.alignment = PP_ALIGN.CENTER

    # Title
    title_box = slide.shapes.add_textbox(Inches(0), Inches(0.1), Inches(13.33), Inches(0.5))
    p = title_box.text_frame.paragraphs[0]
    p.text = "JA BizTown 3.0 — Solution Architecture (RFP Compliant)"
    p.font.bold = True
    p.font.size = Pt(20)
    p.alignment = PP_ALIGN.CENTER

    # Layers
    # Layer 1: Client
    add_layer_bar(0.6, 1.2, "CLIENT & DEVICE LAYER | REACT NATIVE · EXPO", BLUE_LIGHT, BLUE_AZURE)
    client_boxes = [
        ("JA Student Tablet", "React Native + Expo"),
        ("Volunteer Tablet", "Unified POS UI"),
        ("Teacher Tablet", "Monitoring Hub"),
        ("App Distribution", "MDM / Store"),
        ("i18next + XLF", "Offline Discovery"),
        ("IndexedDB", "SQLite Cache"),
        ("Offline Queue", "FIFO Sync"),
        ("Lottie Animation", "Offline Assets")
    ]
    for i, (t, s) in enumerate(client_boxes):
        add_box(0.4 + (i * 1.6), 0.85, t, s, BLUE_AZURE)

    # Layer 2: Security
    add_layer_bar(1.9, 1.2, "PERIMETER SECURITY | ENTRA B2C · KEY VAULT", ORANGE_LIGHT, ORANGE)
    sec_boxes = [
        ("Azure Front Door", "Global Entry"),
        ("WAF (OWASP)", "Threat Filtering"),
        ("Microsoft Entra", "QR/PIN Auth"),
        ("Azure Key Vault", "Secrets / CMK"),
        ("Azure APIM", "Rate Limiting"),
        ("Azure Policy", "RBAC Rules"),
        ("Azure DevOps", "CI/CD Pipeline"),
        ("Azure Monitor", "30-Day Retention")
    ]
    for i, (t, s) in enumerate(sec_boxes):
        add_box(0.4 + (i * 1.6), 2.15, t, s, ORANGE)

    # Layer 3: Compute
    add_layer_bar(3.2, 1.2, "COMPUTE LAYER | AKS · FUNCTIONS · SIGNALR", BLUE_LIGHT, BLUE_AZURE)
    comp_boxes = [
        ("Microsoft AKS", "Microservices"),
        ("Azure Functions", "Serverless Reconciliation"),
        ("Azure SignalR", "Real-time Push"),
        ("Azure Service Bus", "Event Queue"),
        ("Azure Event Hubs", "High-Scale Ingestion"),
        ("Azure CDN", "Media Delivery"),
        ("Azure Logic App", "Workflows")
    ]
    for i, (t, s) in enumerate(comp_boxes):
        add_box(0.4 + (i * 1.8), 3.45, t, s, BLUE_AZURE, RGBColor(225, 245, 254) if "Hub" in t else None)

    # Layer 4: AI
    add_layer_bar(4.5, 1.1, "AI INNOVATION LAYER | GEN AI", PURPLE_LIGHT, PURPLE)
    ai_boxes = [
        ("Azure OpenAI", "GPT-4o Slogans"),
        ("DALL-E 3", "Banner Generation"),
        ("App Insights", "AI Telemetry")
    ]
    for i, (t, s) in enumerate(ai_boxes):
        add_box(4.5 + (i * 1.6), 4.65, t, s, PURPLE)

    # Layer 5: Data
    add_layer_bar(5.7, 1.1, "DATA & CONTENT LAYER | SQL · CMS · BLOB", RED_LIGHT, RED)
    data_boxes = [
        ("Azure SQL", "Simulation Ledger"),
        ("CMS Database", "Headless Backend"),
        ("Azure Blob", "Media Assets"),
        ("Azure Redis", "Live Cache"),
        ("Headless CMS", "Content Layouts"),
        ("SQL Geo-Rep", "DR Support")
    ]
    for i, (t, s) in enumerate(data_boxes):
        add_box(1.0 + (i * 1.8), 5.85, t, s, RED, RGBColor(255, 235, 238) if "CMS" in t else None)

    # Layer 6: Obs
    add_layer_bar(6.9, 0.5, "OBSERVABILITY & COMPLIANCE", GREEN_LIGHT, GREEN)
    obs_boxes = [
        ("Azure Monitor", "Full Audit"),
        ("Log Analytics", "Tracing"),
        ("COPPA Aligned", "Zero Student PII"),
        ("FERPA Aligned", "30-Day Retention")
    ]
    for i, (t, s) in enumerate(obs_boxes):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5 + (i * 1.8)), Inches(7.0), Inches(1.6), Inches(0.35))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
        shape.line.color.rgb = GREEN
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = t
        p.font.size = Pt(8)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    prs.save('Architecture.pptx')
    print("Architecture.pptx created successfully.")

if __name__ == "__main__":
    create_presentation()
